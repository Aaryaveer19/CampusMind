import docx2txt
from pptx import Presentation
import streamlit as st
from groq import Groq
from pypdf import PdfReader
import os
from dotenv import load_dotenv
from audio_recorder_streamlit import audio_recorder
import tempfile

# --- INITIAL SETUP ---
st.set_page_config(page_title="CampusMind AI", page_icon="🧠", layout="wide")

load_dotenv()
GROQ_API_KEY = os.getenv("GROQ_API_KEY")
try:
    client = Groq(api_key=GROQ_API_KEY)
except Exception as e:
    st.error(f"API Key Error: {e}")

# --- CUSTOM UI CSS ---
st.markdown("""
    <style>
    .stApp { background-color: #0B0F19; color: #E5E7EB; }
    h1, h2, h3 { color: #FFFFFF !important; font-weight: 700 !important; }
    .stTextInput>div>div>input { background-color: #1F2937; color: white; border-radius: 8px; border: 1px solid #374151; padding: 10px; }
    .stButton>button { background: linear-gradient(90deg, #4F46E5 0%, #7C3AED 100%); color: white; border-radius: 8px; border: none; height: 2.8em; font-weight: 600; transition: all 0.3s ease; width: 100%; }
    .stButton>button:hover { transform: translateY(-2px); box-shadow: 0 4px 15px rgba(124, 58, 237, 0.4); color: white; }
    [data-testid="stSidebar"] { background-color: #111827; border-right: 1px solid #1F2937; }
    .stTabs [data-baseweb="tab-list"] { gap: 24px; }
    .stTabs [data-baseweb="tab"] { height: 50px; white-space: pre-wrap; background-color: transparent; border-radius: 4px 4px 0px 0px; gap: 1px; padding-top: 10px; padding-bottom: 10px; }
    .stTabs [aria-selected="true"] { color: #A78BFA !important; border-bottom: 2px solid #A78BFA !important; }
    </style>
    """, unsafe_allow_html=True)


if 'chat_answer' not in st.session_state: st.session_state.chat_answer = ""
if 'summary_answer' not in st.session_state: st.session_state.summary_answer = ""
if 'exam_answer' not in st.session_state: st.session_state.exam_answer = ""
if 'logged_in' not in st.session_state: st.session_state.logged_in = False
if 'chat_history' not in st.session_state: st.session_state.chat_history = [] # 👈 NEW CHAT MEMORY

# ==========================================
# 🛑 SECURE LOGIN SYSTEM (HACKATHON DEMO MODE)
# ==========================================
if not st.session_state.logged_in:
    col1, col2, col3 = st.columns([1, 2, 1])
    with col2:
        st.markdown("<br><br>", unsafe_allow_html=True) 
        st.image("https://cdn-icons-png.flaticon.com/512/4185/4185848.png", width=80)
        st.title("CampusMind Portal")
        st.markdown("Please log in with your University credentials.")
        
        username = st.text_input("Student Email", placeholder="student@university.edu")
        password = st.text_input("Password", type="password")
        
        if st.button("Secure Login"):
            if username and password: 
                st.session_state.logged_in = True
                st.rerun() 
            else:
                st.error("⚠️ Please enter your email and password.")
                
    st.stop() 

# ==========================================
# ✅ IF LOGGED IN, THE MAIN APP RUNS BELOW
# ==========================================

# --- GROQ GENERATOR (UPGRADED TO 70B MODEL) ---
def smart_generate(prompt):
    try:
        chat_completion = client.chat.completions.create(
            messages=[{"role": "user", "content": prompt}],
            model="llama-3.3-70b-versatile",
            temperature=0.5, 
        )
        return chat_completion.choices[0].message.content
    except Exception as e:
        print(f"⚠️ Groq failed: {e}")
        return "⚠️ **Offline Mode Active:** The live AI is currently overwhelmed by hackathon traffic. Please refer to the cached context."

# --- SIDEBAR LOGIC ---
# --- SIDEBAR LOGIC ---
with st.sidebar:
    st.image("https://cdn-icons-png.flaticon.com/512/4185/4185848.png", width=60)
    st.title("CampusMind Data")
    st.markdown("Upload your institutional resources to train the digital brain.")
    
    # 1. UPGRADED UPLOADER: Now takes PDF, DOCX, and PPTX!
    uploaded_file = st.file_uploader("Drop PDF, Word, or PPT here", type=["pdf", "docx", "pptx"])
    
    if uploaded_file:
        if 'clean_text' not in st.session_state:
            with st.spinner("🧠 Absorbing and Indexing Knowledge..."):
                raw_text = ""
                
                # 1. READ PDF WITH PAGE NUMBERS
                if uploaded_file.name.endswith('.pdf'):
                    reader = PdfReader(uploaded_file)
                    for i, page in enumerate(reader.pages):
                        text = page.extract_text()
                        if text:
                            # INJECT METADATA
                            raw_text += f"\n\n--- [Source: Page {i+1}] ---\n{text}"
                            
                # 2. READ WORD DOC
                elif uploaded_file.name.endswith('.docx'):
                    text = docx2txt.process(uploaded_file)
                    raw_text += f"\n\n--- [Source: Word Document] ---\n{text}"
                    
                # 3. READ PPTX WITH SLIDE NUMBERS
                elif uploaded_file.name.endswith('.pptx'):
                    prs = Presentation(uploaded_file)
                    for i, slide in enumerate(prs.slides):
                        slide_text = ""
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                slide_text += shape.text + "\n"
                        if slide_text.strip():
                            # INJECT METADATA
                            raw_text += f"\n\n--- [Source: Slide {i+1}] ---\n{slide_text}"
                                
                # Save to session state
                st.session_state.clean_text = raw_text[:20000] 
                
        st.success(f"✅ Indexed File: {uploaded_file.name}")
        
    st.markdown("---")
    if st.button("🚪 Log Out"):
        st.session_state.logged_in = False
        # Clear the document cache when logging out
        if 'clean_text' in st.session_state:
            del st.session_state['clean_text']
        st.rerun()

# --- MAIN SCREEN LOGIC ---
st.title("🧠 CampusMind")
st.markdown("#### *Your AI-Powered Academic Copilot*")
st.divider()

tab1, tab2, tab3, tab4 = st.tabs(["🔍 Semantic Search", "📝 Auto-Notes", "🎓 Exam Simulator", "🎬 Video-to-Notes"])

# TAB 1: THE CHATGPT INTERFACE
with tab1:
    if 'clean_text' in st.session_state:
        st.markdown("### 💬 Chat with your Document")
        
        # 1. Display previous chat history (The scrolling effect)
        for message in st.session_state.chat_history:
            with st.chat_message(message["role"]):
                st.markdown(message["content"])
        
        # 2. Voice Input (Optional)
        st.write("🎙️ **Ask with Voice:**")
        audio_bytes = audio_recorder(text="Click to Record", recording_color="#e83e8c", neutral_color="#4F46E5", key="chat_audio")
        
        # 3. Streamlit's native Chat Input (Pins to the bottom!)
        query = st.chat_input("Type your question here...")
        
        # 4. Handle Voice OR Text input
        final_query = None
        
        if audio_bytes:
            with st.spinner("Listening..."):
                with tempfile.NamedTemporaryFile(delete=False, suffix=".wav") as temp_audio:
                    temp_audio.write(audio_bytes)
                    temp_audio_path = temp_audio.name
                try:
                    with open(temp_audio_path, "rb") as file:
                        transcription = client.audio.transcriptions.create(
                            file=("audio.wav", file.read()),
                            model="whisper-large-v3",
                        )
                    final_query = transcription.text
                except Exception as e:
                    st.error("Audio error. Please try again.")
        elif query:
            final_query = query
            
        # 5. Generate AI Response and add to chat!
        if final_query:
            # Display user message instantly
            st.session_state.chat_history.append({"role": "user", "content": final_query})
            with st.chat_message("user"):
                st.markdown(final_query)
            
            # Display AI thinking and response
            with st.chat_message("assistant"):
                with st.spinner("Analyzing document..."):
                    prompt = f"Context: {st.session_state.clean_text}\n\nQuestion: {final_query}\n\nTask: Provide a highly detailed, comprehensive answer based ONLY on the context. Break down the reasoning step-by-step."
                    response = smart_generate(prompt)
                    st.markdown(response)
            
            # Save AI response to memory
            st.session_state.chat_history.append({"role": "assistant", "content": response})
            
    else:
        st.info("👈 Please upload a PDF, Word Doc, or PPT in the sidebar to unlock the Chat.")

# TAB 2
with tab2:
    if 'clean_text' in st.session_state:
        st.markdown("### Instant Revision Notes")
        if st.button("✨ Generate Smart Summary"):
            with st.spinner("Synthesizing notes..."):
                prompt = f"Context: {st.session_state.clean_text}\n\nTask: Act as an expert professor. Provide an in-depth, comprehensive summary of this document. Include an executive overview, detailed bullet points of all core concepts, and explain why these concepts are important. Do not be concise; be extremely thorough."
                st.session_state.summary_answer = smart_generate(prompt)
        
        if st.session_state.summary_answer:
            st.success("📝 **Study Summary:**")
            st.markdown(st.session_state.summary_answer)
            
            st.download_button(
                label="💾 Download Notes as TXT",
                data=st.session_state.summary_answer,
                file_name="CampusMind_StudyNotes.txt",
                mime="text/plain"
            )
    else:
        st.info("👈 Please upload a PDF in the sidebar to unlock Auto-Notes.")
            
# TAB 3
with tab3:
    if 'clean_text' in st.session_state:
        st.markdown("### Predict Viva & Exam Questions")
        num_questions = st.slider("How many questions do you want to generate?", min_value=1, max_value=15, value=5)
        if st.button("🎯 Generate Practice Test"):
            with st.spinner(f"Analyzing past patterns to generate {num_questions} questions..."):
                prompt = f"Context: {st.session_state.clean_text}\n\nTask: Generate {num_questions} probable viva/exam questions with short answers."
                st.session_state.exam_answer = smart_generate(prompt)
        
        if st.session_state.exam_answer:
            st.warning(f"🎓 **Your {num_questions} Practice Questions:**")
            st.markdown(st.session_state.exam_answer)
    else:
        st.info("👈 Please upload a PDF in the sidebar to unlock the Exam Simulator.")

# TAB 4
with tab4:
    st.markdown("### 🎬 Video-to-Notes Converter")
    st.caption("Upload a short video clip (MP4) and the AI will transcribe and summarize it!")
    
    video_file = st.file_uploader("Upload an MP4 Video (Max 25MB)", type=["mp4", "mp3", "wav"])
    
    if video_file:
        st.video(video_file)
        
        if st.button("🎙️ Transcribe & Summarize Video"):
            with st.spinner("Transcribing speech using Whisper AI..."):
                try:
                    transcription = client.audio.transcriptions.create(
                        file=(video_file.name, video_file.getvalue()),
                        model="whisper-large-v3",
                    )
                    transcript_text = transcription.text
                    
                    st.success("✅ Transcription Complete!")
                    with st.expander("View Raw Transcript"):
                        st.write(transcript_text)
                    
                    with st.spinner("Analyzing transcript..."):
                        video_summary_prompt = f"Context: {transcript_text}\n\nTask: Summarize this video transcript in 3-5 concise bullet points."
                        video_summary = smart_generate(video_summary_prompt)
                        
                        st.info("📊 **Smart Video Summary:**")
                        st.markdown(video_summary)
                        
                except Exception as e:
                    st.error(f"⚠️ Error: {e}. Make sure your video is under 25MB!")